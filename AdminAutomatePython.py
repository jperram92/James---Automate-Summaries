import argparse
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import logging
import sys
import os

# Configure default paths
DEFAULT_INPUT = os.path.join("input", "requirements.xlsx")
DEFAULT_OUTPUT = os.path.join("output", "presentation.pptx")
TEMPLATE_FILE = os.path.join("templates", "General presentation.pptx")

# Set up logging
logging.basicConfig(filename='excel_to_ppt.log', level=logging.INFO,
                    format='%(asctime)s - %(levelname)s - %(message)s')

def validate_input(input_file):
    """Validate input files and structure"""
    if not os.path.exists(input_file):
        raise FileNotFoundError(f"Excel file {input_file} not found")

    required_columns = ['Section', 'Title', 'Description', 'Priority']
    df = pd.read_excel(input_file)
    missing_cols = [col for col in required_columns if col not in df.columns]
    
    if missing_cols:
        raise ValueError(f"Missing required columns: {', '.join(missing_cols)}")
    
    return df

def add_priority_badge(slide, priority):
    """Add colored priority badge to slide"""
    priority_colors = {
        'High': RGBColor(198, 31, 60),    # Red
        'Medium': RGBColor(255, 192, 0),  # Amber
        'Low': RGBColor(59, 168, 85)      # Green
    }
    
    left = Inches(8.5)
    top = Inches(0.2)
    width = Inches(1)
    height = Inches(0.4)
    
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = priority_colors.get(priority, RGBColor(128, 128, 128))
    
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = priority
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = 1  # Center alignment

def add_logo(slide, logo_path):
    """Add logo to slide"""
    if os.path.exists(logo_path):
        left = Inches(0.5)
        top = Inches(0.5)
        height = Inches(1.5)
        slide.shapes.add_picture(logo_path, left, top, height=height)
    else:
        logging.warning(f"Logo not found: {logo_path}")

def set_slide_background(slide, priority):
    """Set a dynamic background color for each slide"""
    background = slide.background
    fill = background.fill
    if priority == 'High':
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 99, 71)  # Light Red for High priority
    elif priority == 'Medium':
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 223, 77)  # Light Yellow for Medium priority
    else:
        fill.solid()
        fill.fore_color.rgb = RGBColor(152, 251, 152)  # Light Green for Low priority

def create_ppt(input_file, output_file, template_file):
    """Main function to create PowerPoint from Excel"""
    try:
        # Create output directory if needed
        os.makedirs(os.path.dirname(output_file), exist_ok=True)
        
        # Load the template presentation
        prs = Presentation(template_file)
        
        # Validate input Excel file
        df = validate_input(input_file)
        
        # Process each row in the Excel input and create slides
        for _, row in df.iterrows():
            slide = prs.slides.add_slide(prs.slide_layouts[2])  # Update layout index based on printout
            
            # Set background color based on priority
            set_slide_background(slide, row['Priority'])
            
            # Set title and content
            title = slide.shapes.title
            title.text = f"{row['Section']} - {row['Title']}"
            title.text_frame.paragraphs[0].font.size = Pt(20)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black title color
            
            # Set content with improved text styling
            content = slide.shapes.placeholders[1]  # Use the correct placeholder index based on layout
            content.text = row['Description']
            content.text_frame.paragraphs[0].font.size = Pt(14)
            content.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black content color
            
            # Add diagram if specified
            if 'Diagram Needed' in row and pd.notna(row['Diagram Needed']):
                img_path = os.path.join("diagrams", f"{row['Diagram Needed']}.png")
                if os.path.exists(img_path):
                    left = Inches(1)
                    top = Inches(2)
                    height = Inches(4)
                    slide.shapes.add_picture(img_path, left, top, height=height)
                else:
                    logging.warning(f"Diagram not found: {img_path}")
            
            # Add priority badge
            add_priority_badge(slide, row['Priority'])
            
            # Add logo if specified
            if 'Logo' in row and pd.notna(row['Logo']):
                add_logo(slide, row['Logo'])
        
        prs.save(output_file)
        print(f"Successfully created {output_file} with {len(df)} slides")
        logging.info(f"Created presentation: {output_file}")
        
    except Exception as e:
        logging.error(f"Error processing file: {str(e)}")
        sys.exit(f"Error: {str(e)}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Convert Excel to PowerPoint')
    parser.add_argument('--input', 
                        default=DEFAULT_INPUT,
                        help=f'Input Excel file (default: {DEFAULT_INPUT})')
    parser.add_argument('--output', 
                        default=DEFAULT_OUTPUT,
                        help=f'Output PPTX file (default: {DEFAULT_OUTPUT})')
    parser.add_argument('--template', 
                        default=TEMPLATE_FILE,
                        help=f'Template PowerPoint file (default: {TEMPLATE_FILE})')
    
    args = parser.parse_args()
    
    create_ppt(
        input_file=args.input,
        output_file=args.output,
        template_file=args.template
    )
